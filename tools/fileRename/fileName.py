import os
from docx import Document
from openpyxl import load_workbook
from pptx import Presentation
from pdfminer.high_level import extract_text
import tkinter as tk
from tkinter import filedialog, messagebox
from tkinter import ttk

def get_first_line_from_word(file_path):
    try:
        doc = Document(file_path)
        return doc.paragraphs[0].text if doc.paragraphs else ''
    except Exception as e:
        print(f"Error reading Word file {file_path}: {e}")
        return ''

def get_first_line_from_excel(file_path):
    try:
        wb = load_workbook(file_path)
        sheet = wb.active
        for row in sheet.iter_rows(min_row=1, values_only=True):
            for value in row:
                if value is not None:
                    return str(value)
        return ''
    except Exception as e:
        print(f"Error reading Excel file {file_path}: {e}")
        return ''

def get_first_line_from_ppt(file_path):
    try:
        presentation = Presentation(file_path)
        for slide in presentation.slides:
            for shape in slide.shapes:
                if hasattr(shape, 'text') and shape.text:
                    return shape.text[:shape.text.index('\n')] if '\n' in shape.text else shape.text
        return ''
    except Exception as e:
        print(f"Error reading PPT file {file_path}: {e}")
        return ''

def get_first_line_from_pdf(file_path):
    try:
        text = extract_text(file_path)
        return text.split('\n', 1)[0] if text else ''
    except Exception as e:
        print(f"Error reading PDF file {file_path}: {e}")
        return ''

def rename_files(directory):
    for filename in os.listdir(directory):
        if filename.lower().endswith(('.docx', '.xlsx', '.pptx', '.pdf')) and not filename.lower().startswith('.~') :
            file_path = os.path.join(directory, filename)
            first_line = ''
            if filename.lower().endswith('.docx'):
                first_line = get_first_line_from_word(file_path)
            elif filename.lower().endswith('.xlsx'):
                first_line = get_first_line_from_excel(file_path)
            elif filename.lower().endswith('.pptx'):
                first_line = get_first_line_from_ppt(file_path)
            elif filename.lower().endswith('.pdf'):
                first_line = get_first_line_from_pdf(file_path)
            
            if first_line:
                new_filename = first_line.strip() + os.path.splitext(filename)[1]
                new_file_path = os.path.join(directory, new_filename)
                os.rename(file_path, new_file_path)
                print(f"Renamed {filename} to {new_filename}")
            else:
                print(f"No first line found for {filename}")


# Specify the directory containing the files
directory = '/Users/zhan/Documents/test'  # Replace with the path to your directory      


# 这里是之前定义的rename_files函数和子函数
def choose_directory():
    directory = filedialog.askdirectory()
    if directory:
        entry.delete(0, tk.END)
        entry.insert(0, directory)

def rename_files_with_ui():
    directory = entry.get()
    if not directory:
        messagebox.showerror("错误", "请选择一个文件夹")
        return
    if not os.path.isdir(directory):
        messagebox.showerror("错误", "所选路径不是一个文件夹")
        return
    try:
        rename_files(directory)
        messagebox.showinfo("完成", "文件重命名完成")
    except Exception as e:
        messagebox.showerror("错误", f"发生错误: {e}")


# 创建主窗口
root = tk.Tk()
root.title("文件批量重命名工具")

# 创建一个标签和输入框用于显示选择的文件夹路径
label = ttk.Label(root, text="请选择文件夹:")
label.pack()
entry = tk.Entry(root, width=30)
entry.pack()

# 创建一个按钮，点击时弹出选择文件夹的对话框
browse_button = ttk.Button(root, text="浏览", command=choose_directory)
browse_button.pack()

# 创建一个按钮，点击时执行重命名操作
rename_button = ttk.Button(root, text="确定", command=rename_files_with_ui)
rename_button.pack()

# 设置窗口的尺寸
width = 350  # 宽度
height = 200  # 高度

# 获取屏幕的宽度和高度
screen_width = root.winfo_screenwidth()
screen_height = root.winfo_screenheight()

# 计算窗口的中心坐标
x = (screen_width / 2) - (width / 2)
y = (screen_height / 2) - (height / 2)

# 将窗口放置在屏幕中心
root.geometry(f'{width}x{height}+{int(x)}+{int(y)}')

# 运行主循环
root.mainloop()

